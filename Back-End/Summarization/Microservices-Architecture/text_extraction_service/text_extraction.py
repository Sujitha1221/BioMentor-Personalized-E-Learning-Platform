import os
import re
import logging
from tqdm import tqdm
from docx import Document
from pptx import Presentation
from autocorrect import Speller
import fitz  # PyMuPDF
import tabula
from PIL import Image
import numpy as np
import easyocr

# Configure logging
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")

# Initialize the spell checker
spell = Speller(lang="en")


def extract_text_and_tables_from_word(docx_path):
    """
    Extract text and tables from Word (.docx) documents.
    """
    try:
        text = ""
        doc = Document(docx_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text.strip() + "\n"
        for table in doc.tables:
            text += convert_table_to_paragraph(table) + "\n"
        logging.info(f"Extracted content from Word document: {docx_path}")
        return text
    except Exception as e:
        logging.error(f"Failed to extract content from Word document: {e}")
        raise RuntimeError(f"Error processing Word document: {docx_path}")


def extract_text_and_tables_from_pptx(pptx_path):
    """
    Extract text and tables from PowerPoint (.pptx) slides.
    """
    try:
        text = ""
        presentation = Presentation(pptx_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            text += paragraph.text.strip() + "\n"
                if shape.has_table:
                    table = shape.table
                    text += convert_table_to_paragraph_from_pptx(table) + "\n"
        logging.info(f"Extracted content from PowerPoint presentation: {pptx_path}")
        return text
    except Exception as e:
        logging.error(f"Failed to extract content from PowerPoint presentation: {e}")
        raise RuntimeError(f"Error processing PowerPoint presentation: {pptx_path}")


def extract_text_from_txt(txt_path):
    """
    Extract text from plain text (.txt) files.
    """
    try:
        with open(txt_path, 'r', encoding='utf-8') as file:
            text = file.read()
        logging.info(f"Extracted text from TXT file: {txt_path}")
        return text
    except Exception as e:
        logging.error(f"Failed to extract text from TXT file: {e}")
        raise RuntimeError(f"Error processing TXT file: {txt_path}")


def convert_table_to_paragraph(table):
    """
    Convert Word table content into readable paragraph text.
    """
    try:
        sentences = []
        headers = [cell.text.strip() for cell in table.rows[0].cells]
        for row in table.rows[1:]:
            row_data = [cell.text.strip() for cell in row.cells]
            if len(headers) == len(row_data):
                sentence = ", ".join(f"{headers[i]}: {row_data[i]}" for i in range(len(headers)))
                sentences.append(sentence)
        return "\n".join(sentences)
    except Exception as e:
        logging.error(f"Failed to convert Word table to paragraph: {e}")
        return f"Error processing table."


def convert_table_to_paragraph_from_pptx(table):
    """
    Convert PowerPoint table content into readable paragraph text.
    """
    try:
        sentences = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            if len(row_data) > 0:
                sentence = ", ".join(row_data)
                sentences.append(sentence)
        return "\n".join(sentences)
    except Exception as e:
        logging.error(f"Failed to convert PowerPoint table to paragraph: {e}")
        return f"Error processing table."


def clean_text(text):
    """
    Clean and structure text by fixing word splits, removing unwanted spaces,
    correcting grammar, and ensuring readability.
    """
    try:
        # 1. Remove headings and chapter titles
        text = re.sub(r'(?i)\b(chapter\s*\d+|section\s*\d+|introduction|table of contents)\b.*', '', text, flags=re.IGNORECASE)
        text = re.sub(r'^[A-Z\s]{3,}\n', '', text, flags=re.MULTILINE)

        # 2. Remove bullet points, list numbers, and standalone numbers
        text = re.sub(r'^[\-\•\*\+\>\➢]+\s*', '', text, flags=re.MULTILINE)  # Bullet points
        text = re.sub(r'^\d+\.\s*', '', text, flags=re.MULTILINE)  # List numbers
        text = re.sub(r'\b\d+\b', '', text)  # Standalone numbers

        # 3. Fix improper word splits (e.g., "m ean" -> "mean")
        text = re.sub(r'(\w)\s+(\w)', r'\1 \2', text)  # Ensure proper spacing between letters
        text = re.sub(r'(\w)\s{2,}(\w)', r'\1 \2', text)  # Remove extra spaces between words

        # 4. Normalize spaces and line breaks
        text = re.sub(r'\s+', ' ', text).strip()
        text = re.sub(r'\n+', '\n', text).strip()

        # 5. Remove redundant phrases and repetitive content
        text = re.sub(r'\b(the study of living organisms is called biology|biology is called biology)\b', 'biology', text, flags=re.IGNORECASE)

        # 6. Fix grammar, punctuation, and sentence structure with progress bar
        sentences = re.split(r'(?<=[.!?])\s+', text)  # Split into sentences
        corrected_sentences = []

        for sentence in tqdm(sentences, desc="Cleaning text sentences"):
            sentence = spell(sentence.strip())  # Fix spelling
            if len(sentence) > 1:
                corrected_sentences.append(sentence[0].upper() + sentence[1:].lower())

        cleaned_text = ' '.join(corrected_sentences).strip()

        # Ensure proper punctuation at the end of the text
        if cleaned_text and cleaned_text[-1] not in '.!?':
            cleaned_text += '.'

        logging.info("Cleaned and formatted text successfully.")
        return cleaned_text
    except Exception as e:
        logging.error(f"Failed to clean text: {e}")
        raise RuntimeError("Error cleaning text.")



def format_as_paragraph(text):
    """
    Format cleaned text into cohesive paragraphs.
    """
    try:
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        paragraph = " ".join(lines)
        paragraph = re.sub(r'\s+', ' ', paragraph).strip()
        logging.info("Formatted text into paragraphs successfully.")
        return paragraph
    except Exception as e:
        logging.error(f"Failed to format text as a paragraph: {e}")
        raise RuntimeError("Error formatting text.")

# Error detection in text
def find_errors(text):
    """
    Identify common errors in the extracted text.
    """
    errors = []

    # Find improperly split words (e.g., 'm ean', 'env ironman')
    split_words = re.findall(r'\b\w\s+\w\b', text)
    if split_words:
        errors.append(f"Improperly split words: {', '.join(split_words)}")

    # Find redundant spaces
    redundant_spaces = re.findall(r'\s{2,}', text)
    if redundant_spaces:
        errors.append("Redundant spaces found in text.")

    # Check for grammar issues: Repeated phrases or words
    repeated_phrases = re.findall(r'\b(\w+ \w+)\b(?=.*\b\1\b)', text)
    if repeated_phrases:
        errors.append(f"Repeated phrases: {', '.join(set(repeated_phrases))}")

    # Identify missing punctuation at sentence endings
    sentences = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s', text)
    for idx, sentence in enumerate(sentences):
        if not sentence.endswith(('.', '!', '?')):
            errors.append(f"Sentence missing punctuation: {sentence[:50]}...")

    # Detect formatting issues (e.g., double punctuation)
    double_punctuation = re.findall(r'[.!?,]{2,}', text)
    if double_punctuation:
        errors.append(f"Double punctuation found: {', '.join(double_punctuation)}")

    # Find invalid words or placeholders (e.g., 'dis ease')
    invalid_words = re.findall(r'\bdis ease\b', text)
    if invalid_words:
        errors.append(f"Invalid words: {', '.join(invalid_words)}")

    # Summarize errors
    if not errors:
        return "No errors found."
    return "\n".join(errors)

def resolve_errors(text):
    """
    Resolve common errors in the extracted text.
    """
    try:
        # 1. Fix improperly split words (e.g., 'i s' -> 'is', 'o r' -> 'or')
        text = re.sub(r'\b(i s|o r|i t|i t s)\b', lambda m: m.group(0).replace(" ", ""), text)

        # 2. Remove repeated phrases (e.g., "the science of the science of")
        repeated_phrases = re.findall(r'\b(\w+\s+\w+)\b(?=.*\b\1\b)', text)
        for phrase in set(repeated_phrases):
            text = re.sub(rf'\b({phrase})\b(?=.*\b\1\b)', phrase, text)

        # 3. Fix double punctuation (e.g., '... -> .')
        text = re.sub(r'[.!?]{2,}', '.', text)

        # 4. Fix invalid words (e.g., 'dis ease' -> 'disease')
        text = re.sub(r'\bdis ease\b', 'disease', text)

        # 5. Normalize spacing and capitalization
        text = re.sub(r'\s+', ' ', text).strip()  # Normalize spaces
        text = re.sub(r'\s+([.,!?])', r'\1', text)  # Remove spaces before punctuation
        sentences = re.split(r'(?<=[.!?])\s+', text)  # Split into sentences
        resolved_sentences = [
            sentence[0].upper() + sentence[1:] if sentence else ""
            for sentence in sentences
        ]
        resolved_text = ' '.join(resolved_sentences).strip()

        # Ensure proper punctuation at the end
        if resolved_text and resolved_text[-1] not in '.!?':
            resolved_text += '.'

        logging.info("Resolved text errors successfully.")
        return resolved_text
    except Exception as e:
        logging.error(f"Error resolving text: {e}")
        raise RuntimeError("Error resolving text.")
    
def clean_table(df):
    """
    Clean and structure the table data from a DataFrame.
    """
    try:
        df = df.dropna(how='all')  # Drop rows with all NaN values
        df = df.dropna(axis=1, how='all')  # Drop columns with all NaN values
        df = df.reset_index(drop=True)  # Reset index after dropping
        df.columns = [col.strip() for col in df.columns]  # Strip whitespace from column names
        df = df[~df.apply(lambda row: row.astype(str).str.contains('©|Rights Reserved').any(), axis=1)]  # Remove noisy rows
        return df
    except Exception as e:
        logging.error(f"Failed to clean table: {e}")
        raise



def extract_text_from_pdf(pdf_path):
    """
    Extract text from PDF using PyMuPDF (fitz) to handle text-based PDFs with proper spacing.
    """
    try:
        text = ""
        document = fitz.open(pdf_path)
        for page_number in range(len(document)):
            page = document[page_number]
            page_text = page.get_text("text")  # Extract text with proper spacing
            if page_text.strip():
                text += page_text.strip() + "\n"
        document.close()
        if not text.strip():
            logging.warning(f"No extractable text found in PDF: {pdf_path}")
        return text
    except Exception as e:
        logging.error(f"Failed to extract text from PDF: {e}")
        raise RuntimeError(f"Error processing PDF: {pdf_path}")


def extract_text_from_image_pdf(pdf_path):
    """
    Extract text from image-based PDFs using fitz (PyMuPDF) with OCR.
    """
    try:
        text = ""
        document = fitz.open(pdf_path)
        reader = easyocr.Reader(["en"], gpu=False)  # Initialize OCR

        for page_number in range(len(document)):
            page = document[page_number]

            # Try direct text extraction
            page_text = page.get_text()
            if page_text.strip():
                text += page_text + "\n"
                continue

            # Extract images from the page
            images = page.get_images(full=True)
            if images:
                for img_index, img in enumerate(images):
                    xref = img[0]
                    pix = fitz.Pixmap(document, xref)
                    if pix.n < 5:  # GRAY or RGB
                        pix_image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    else:  # CMYK: Convert to RGB
                        pix_image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples, "raw", pix.colorspace, 0, 1)

                    # Convert image to array for OCR
                    img_array = np.array(pix_image)
                    ocr_result = reader.readtext(img_array, detail=0)  # Extract text
                    page_ocr_text = " ".join(ocr_result).strip()

                    if page_ocr_text:
                        text += page_ocr_text + "\n"
                        logging.info(f"Extracted text from image on page {page_number + 1}, image {img_index + 1}.")

        document.close()

        if not text.strip():
            logging.warning(f"No extractable text or embedded data found in image-based PDF: {pdf_path}")
        else:
            logging.info(f"Processed image-based PDF using fitz and OCR: {pdf_path}")

        return text.strip()
    except Exception as e:
        logging.error(f"Failed to process image-based PDF with fitz and OCR: {e}")
        raise RuntimeError(f"Error processing image-based PDF: {e}")


# def extract_text_from_image_pdf(pdf_path):
#     """
#     Optimized version: Extract text from image-based PDFs using fitz (PyMuPDF) and only perform OCR
#     when necessary to minimize processing time.
#     """
#     try:
#         text = ""
#         document = fitz.open(pdf_path)

#         # Initialize EasyOCR Reader once
#         reader = easyocr.Reader(["en"], gpu=False)  # Use GPU=True if available for faster OCR

#         for page_number in range(len(document)):
#             page = document[page_number]

#             # Try direct text extraction first
#             page_text = page.get_text()
#             if page_text.strip():  # If text is found, skip OCR
#                 text += page_text + "\n"
#                 continue

#             # If no text is found, perform OCR
#             pix = page.get_pixmap()  # Render page as image
#             img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
#             ocr_result = reader.readtext(np.array(img), detail=0)  # Perform OCR
#             ocr_text = " ".join(ocr_result).strip()

#             if ocr_text:
#                 text += ocr_text + "\n"
#                 logging.info(f"Extracted text from page {page_number + 1} using OCR.")

#         document.close()

#         if not text.strip():
#             logging.warning(f"No extractable text or embedded data found in PDF: {pdf_path}")
#         else:
#             logging.info(f"Processed PDF successfully: {pdf_path}")

#         return text.strip()
#     except Exception as e:
#         logging.error(f"Failed to process PDF: {e}")
#         raise RuntimeError(f"Error processing PDF: {e}")
    
def extract_tables_from_pdf(pdf_path):
    """
    Extract and clean tables from PDF using Tabula.
    """
    try:
        tables_text = ""
        tables = tabula.read_pdf(pdf_path, pages="all", multiple_tables=True, silent=True)
        for table_index, table in enumerate(tables):
            cleaned_table = clean_table(table)  # Clean the table data
            if not cleaned_table.empty:
                tables_text += f"Table {table_index + 1}:\n{cleaned_table.to_string(index=False)}\n\n"
        return tables_text
    except Exception as e:
        logging.error(f"Failed to extract tables from PDF: {e}")
        raise RuntimeError(f"Error processing tables from PDF: {pdf_path}")
    
def extract_content_from_pdf(pdf_path):
    """
    Extract text and tables from a PDF, using both text-based and image-based methods.
    """
    try:
        logging.info(f"Extracting content from PDF: {pdf_path}")

        # Extract text using PyMuPDF
        extracted_text = extract_text_from_pdf(pdf_path)

        # If no text is found, attempt OCR for image-based PDFs
        if not extracted_text.strip():
            logging.info("No text found in PDF. Attempting OCR on image-based PDF...")
            extracted_text = extract_text_from_image_pdf(pdf_path)

        # Extract tables and append to content
        extracted_tables = extract_tables_from_pdf(pdf_path)
        if extracted_tables.strip():
            extracted_text += "\n\n" + extracted_tables

        logging.info(f"Extraction complete for PDF: {pdf_path}")
        return extracted_text.strip()
    except Exception as e:
        logging.error(f"Failed to extract content from PDF: {e}")
        raise RuntimeError(f"Error processing PDF: {pdf_path}")


def extract_content(file_path):
    """
    Extract text and tables from a given file, clean, and format it.
    """
    try:
        if file_path.endswith(".docx"):
            raw_text = extract_text_and_tables_from_word(file_path)
        elif file_path.endswith(".pptx"):
            raw_text = extract_text_and_tables_from_pptx(file_path)
        elif file_path.endswith(".pdf"):
            raw_text = extract_content_from_pdf(file_path)
        elif file_path.endswith(".txt"):
            raw_text = extract_text_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_path}")

        # Clean and format the extracted raw text
        cleaned_text = clean_text(raw_text)

        # Optionally format as paragraphs
        formatted_text = format_as_paragraph(cleaned_text)

        # Find errors in the cleaned text
        errors = find_errors(formatted_text)
        print("\nErrors Found:\n", errors)

        # # Resolve errors
        resolved_text = resolve_errors(formatted_text)
        print("Cleaned text:", resolved_text)
        return formatted_text
    except Exception as e:
        logging.error(f"Failed to process file: {e}")
        raise RuntimeError(f"Error processing file: {file_path}")


